#!/usr/bin/env python3
"""
Yuno Business Case Generator — Template Editor (bridge-adapted)

Reads config JSON from stdin, writes generated PPTX bytes to stdout.
Used by the bridge's /api/generate-business-case endpoint.

Also supports CLI mode: python generate.py config.json [template.pptx] [output.pptx]
"""
import json
import sys
import os
import subprocess
import io
import re
import tempfile
from pathlib import Path
from copy import deepcopy

try:
    from pptx import Presentation
    from pptx.util import Emu, Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


# ============ FORMATTERS ============
def fmt_k(n):
    if n >= 1_000_000:
        return f"${n/1_000_000:.1f}M"
    if n >= 1000:
        return f"${n/1000:.0f}K"
    return f"${n:,.0f}"


def fmt_money(n):
    return f"${n:,.0f}"


# ============ TEXT REPLACEMENT HELPERS ============
def replace_in_shape(shape, old, new):
    if not shape.has_text_frame:
        return 0
    count = 0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                count += 1
        full_text = "".join(r.text for r in para.runs)
        if old in full_text and not any(old in r.text for r in para.runs):
            new_text = full_text.replace(old, new)
            if para.runs:
                para.runs[0].text = new_text
                for r in list(para.runs[1:]):
                    r_elem = r._r
                    r_elem.getparent().remove(r_elem)
                count += 1
    return count


def set_cell_text_preserve(cell, new_text):
    tf = cell.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = new_text
        for run in list(tf.paragraphs[0].runs[1:]):
            r_elem = run._r
            r_elem.getparent().remove(r_elem)
        for p in list(tf.paragraphs[1:]):
            p_elem = p._p
            p_elem.getparent().remove(p_elem)


def set_cell_text_multiline(cell, lines):
    tf = cell.text_frame
    if not tf.paragraphs:
        cell.text = "\n".join(lines)
        return
    if tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = lines[0]
        for run in list(tf.paragraphs[0].runs[1:]):
            r_elem = run._r
            r_elem.getparent().remove(r_elem)
    if len(lines) > 1:
        for i, line in enumerate(lines[1:], 1):
            if i < len(tf.paragraphs):
                p = tf.paragraphs[i]
                if p.runs:
                    p.runs[0].text = line
                    for run in list(p.runs[1:]):
                        r_elem = run._r
                        r_elem.getparent().remove(r_elem)
    for p in list(tf.paragraphs[len(lines):]):
        p_elem = p._p
        p_elem.getparent().remove(p_elem)


# ============ LOGO FALLBACK ============
def create_logo_fallback(company_name, color_hex, output_path):
    """Create a simple branded logo using PIL (avoids Node/sharp dependency)."""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        # Create a simple colored square as last resort
        with open(output_path, "wb") as f:
            # 1x1 white PNG
            f.write(bytes.fromhex("89504e470d0a1a0a0000000d4948445200000001000000010100000000376ef9240000000a49444154789c6300010000000500010d0a2db40000000049454e44ae426082"))
        return

    size = 600
    img = Image.new("RGB", (size, size), "white")
    draw = ImageDraw.Draw(img)
    r, g, b = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)
    # Try common fonts
    font = None
    for font_path in ["/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", "/System/Library/Fonts/Helvetica.ttc", "Arial.ttf"]:
        try:
            font_size = min(120, max(40, int(size / max(1, len(company_name)) * 1.5)))
            font = ImageFont.truetype(font_path, font_size)
            break
        except:
            continue
    if font is None:
        font = ImageFont.load_default()
    # Center text
    text = company_name.upper()
    bbox = draw.textbbox((0, 0), text, font=font)
    w, h = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text(((size - w) / 2, (size - h) / 2 - bbox[1]), text, fill=(r, g, b), font=font)
    img.save(output_path, "PNG")


# ============ MAIN TEMPLATE EDITOR ============
def edit_template(config, template_path, output_path, logo_path=None):
    prs = Presentation(template_path)

    company = config["clientName"]
    mdr_actual = config["mdrActual"]
    mdr_nuevo = config["mdrNuevo"]
    total_txns = config["totalTxnMes"]
    ticket = config["ticketPromedio"]
    tpv_mensual = total_txns * ticket
    costo_actual = tpv_mensual * mdr_actual
    costo_nuevo = tpv_mensual * mdr_nuevo
    reduccion_mdr = costo_actual - costo_nuevo

    delta_aprobacion = config.get("deltaAprobacion", 0.03)
    aumento_tpv = tpv_mensual * delta_aprobacion
    margen = config.get("margenProducto", 0.30)
    revenue_aumento = aumento_tpv * margen

    ahorro_conciliacion = config.get("ahorroConciliacion", 5000)
    ahorro_operativo = config.get("ahorroOperativo", 5000)
    total_impacto = revenue_aumento + reduccion_mdr + ahorro_conciliacion + ahorro_operativo

    countries = config["countries"]
    pricing_type = config.get("pricingType", "flat")
    flat_price = config.get("flatPrice", 0.10)
    tramos = config.get("tramos", [])
    min_transaccional = config.get("minimoTransaccional", "No hay mínimo transaccional")
    saas_fee = config.get("saasFee", "")

    bdm_name = config.get("bdmName", "Rasheed Bayter")
    bdm_phone = config.get("bdmPhone", "+57 3164730919")
    bdm_email = config.get("bdmEmail", "rasheed@y.uno")

    mdr_antes_str = f"{mdr_actual*100:.1f}%"
    mdr_despues_str = f"{mdr_nuevo*100:.1f}%"
    reduccion_mdr_str = fmt_money(reduccion_mdr)
    aumento_tpv_str = fmt_money(aumento_tpv)
    aumento_aprobacion_str = f"{int(delta_aprobacion*100)}%"

    print(f"Editing for {company}: TPV={fmt_money(tpv_mensual)}, MDR {mdr_antes_str}→{mdr_despues_str}, Total={fmt_k(total_impacto)}", file=sys.stderr)

    # SLIDE 1: Logo
    slide1 = prs.slides[0]
    if logo_path and os.path.exists(logo_path):
        logo_shape = None
        for shape in slide1.shapes:
            if shape.has_text_frame and "{{LOGO}}" in shape.text_frame.text:
                logo_shape = shape
                break
        if logo_shape:
            left, top, height = logo_shape.left, logo_shape.top, logo_shape.height
            sp = logo_shape._element
            sp.getparent().remove(sp)
            slide1.shapes.add_picture(logo_path, left, top, width=height, height=height)

    # SLIDE 4: Company name
    if len(prs.slides) > 3:
        for shape in prs.slides[3].shapes:
            replace_in_shape(shape, "{{COMPANY_NAME}}", company)

    # SLIDE 5: Two logos
    if len(prs.slides) > 4 and logo_path and os.path.exists(logo_path):
        slide5 = prs.slides[4]
        logo_shapes = [s for s in slide5.shapes if s.has_text_frame and "{{logo}}" in s.text_frame.text]
        for lshape in logo_shapes:
            left, top, width, height = lshape.left, lshape.top, lshape.width, lshape.height
            sp = lshape._element
            sp.getparent().remove(sp)
            logo_size = height
            new_left = left + (width - logo_size) // 2
            slide5.shapes.add_picture(logo_path, new_left, top, width=logo_size, height=logo_size)

    # SLIDE 8: Optimization text
    if len(prs.slides) > 7:
        replacements = [
            ("{{COMPANY_NAME}}", company),
            ("{{costo antes}}", mdr_antes_str),
            ("{{costo despues}}", mdr_despues_str),
            ("{{reduccion_mdr}}", reduccion_mdr_str),
            ("{{aumento_aprobacion}}", aumento_aprobacion_str),
            ("{{aumento_tpv}}", aumento_tpv_str),
        ]
        for shape in prs.slides[7].shapes:
            for old, new in replacements:
                replace_in_shape(shape, old, new)

    # SLIDE 10: Volume table
    if len(prs.slides) > 9:
        slide10 = prs.slides[9]
        table_s10 = None
        for shape in slide10.shapes:
            if shape.has_table:
                table_s10 = shape.table
                break
        if table_s10 is not None:
            tbl = table_s10._tbl
            ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            trs = tbl.findall(f".//{ns}tr")
            num_countries = len(countries)
            # Remove rows num_countries+1..14 (keep header, data, total)
            rows_to_remove = list(range(num_countries + 1, 15))
            for idx in sorted(rows_to_remove, reverse=True):
                if idx < len(trs):
                    tr = trs[idx]
                    tr.getparent().remove(tr)
            for shape in slide10.shapes:
                if shape.has_table:
                    table_s10 = shape.table
                    break
            total_txns_sum = 0
            total_tpv_mes = 0
            for i, country in enumerate(countries):
                if i + 1 >= len(table_s10.rows):
                    break
                row = table_s10.rows[i + 1]
                txns = country.get("txnsPerMonth", country.get("txnPerMonth", 0))
                ct = country.get("ticketPromedio", ticket)
                mdr = country.get("mdr", mdr_actual)
                tpv_mes = txns * ct
                tpv_year = tpv_mes * 12
                costo = tpv_mes * mdr
                set_cell_text_preserve(row.cells[0], country.get("name", country.get("country", "?")))
                set_cell_text_preserve(row.cells[1], f"${tpv_year:,.0f}")
                set_cell_text_preserve(row.cells[2], f"{txns:,}")
                set_cell_text_preserve(row.cells[3], f"{mdr*100:.1f}%")
                set_cell_text_preserve(row.cells[4], f"${costo:,.0f}")
                total_txns_sum += txns
                total_tpv_mes += tpv_mes
            num_rows = len(table_s10.rows)
            total_row = table_s10.rows[num_rows - 1]
            total_tpv_year = total_tpv_mes * 12
            total_costo = total_tpv_mes * mdr_actual
            set_cell_text_preserve(total_row.cells[0], "Total")
            set_cell_text_preserve(total_row.cells[1], f"${total_tpv_year:,.0f}")
            set_cell_text_preserve(total_row.cells[2], f"{total_txns_sum:,}")
            set_cell_text_preserve(total_row.cells[3], mdr_antes_str)
            set_cell_text_preserve(total_row.cells[4], f"${total_costo:,.0f}")

    # SLIDE 12: Efecto Yuno
    if len(prs.slides) > 11:
        slide12 = prs.slides[11]
        table_s12 = None
        total_shape_s12 = None
        for shape in slide12.shapes:
            if shape.has_table:
                table_s12 = shape.table
            elif shape.has_text_frame and "Total Impacto" in shape.text_frame.text:
                total_shape_s12 = shape
        if table_s12 is not None:
            margen_pct = int(margen * 100)
            delta_pct = int(delta_aprobacion * 100)
            set_cell_text_multiline(table_s12.rows[1].cells[0], [
                "Aumento de Aprobación",
                f"({margen_pct}% margen sobre TPV)"
            ])
            set_cell_text_multiline(table_s12.rows[1].cells[1], [
                f"+{delta_pct}% = {fmt_k(aumento_tpv)} TPV",
                f"× {margen_pct}% margen"
            ])
            set_cell_text_preserve(table_s12.rows[1].cells[2], fmt_k(revenue_aumento))
            set_cell_text_preserve(table_s12.rows[2].cells[0], "Reducción Comisión (MDR)")
            set_cell_text_preserve(table_s12.rows[2].cells[1], f"De {mdr_antes_str} a {mdr_despues_str}")
            set_cell_text_preserve(table_s12.rows[2].cells[2], fmt_k(reduccion_mdr))
            set_cell_text_preserve(table_s12.rows[3].cells[2], fmt_money(ahorro_conciliacion))
            set_cell_text_preserve(table_s12.rows[4].cells[2], fmt_money(ahorro_operativo))
        if total_shape_s12:
            for para in total_shape_s12.text_frame.paragraphs:
                full = "".join(r.text for r in para.runs)
                if "Total Impacto" in full:
                    new_full = re.sub(r"\$[\d,]+K?", fmt_k(total_impacto), full)
                    if new_full != full and para.runs:
                        para.runs[0].text = new_full
                        for r in list(para.runs[1:]):
                            r_elem = r._r
                            r_elem.getparent().remove(r_elem)

    # SLIDE 13: Company name (template has typo)
    if len(prs.slides) > 12:
        for shape in prs.slides[12].shapes:
            replace_in_shape(shape, "{{COMPANY_NAME}", company)
            replace_in_shape(shape, "{{COMPANY_NAME}}", company)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if f"{company}podría" in run.text:
                            run.text = run.text.replace(f"{company}podría", f"{company} podría")

    # SLIDE 14: Impact summary
    if len(prs.slides) > 13:
        shapes_s14 = list(prs.slides[13].shapes)
        s14_updates = {7: fmt_k(total_impacto), 11: fmt_k(reduccion_mdr), 13: fmt_k(revenue_aumento)}
        for idx, new_val in s14_updates.items():
            if idx < len(shapes_s14):
                shape = shapes_s14[idx]
                if shape.has_text_frame and shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    if para.runs:
                        para.runs[0].text = new_val

    # SLIDE 16: Pricing
    if len(prs.slides) > 15:
        slide16 = prs.slides[15]
        for shape in slide16.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if "Mínimo transaccional" in text and "20,000 transacciones/mes" in text:
                    replace_in_shape(shape, "20,000 transacciones/mes", min_transaccional)
                elif "Fee fijo SaaS mensual" in text:
                    if saas_fee:
                        replace_in_shape(shape, "5,000 USD", saas_fee)
                    else:
                        for para in shape.text_frame.paragraphs:
                            full = "".join(r.text for r in para.runs)
                            if "Fee fijo SaaS" in full:
                                new_full = full.split(".")[0] + ". Sin fee fijo SaaS mensual"
                                if para.runs:
                                    para.runs[0].text = new_full
                                    for r in list(para.runs[1:]):
                                        r_elem = r._r
                                        r_elem.getparent().remove(r_elem)
        table_s16 = None
        discount_shape = None
        for shape in slide16.shapes:
            if shape.has_table:
                table_s16 = shape
            elif shape.has_text_frame and "descuento progresivo" in shape.text_frame.text:
                discount_shape = shape
        if pricing_type == "flat" and table_s16 is not None:
            tb_left, tb_top, tb_width = table_s16.left, table_s16.top, table_s16.width
            sp = table_s16._element
            sp.getparent().remove(sp)
            tb = slide16.shapes.add_textbox(tb_left, tb_top + Inches(0.5), tb_width, Inches(1.5))
            tb.text_frame.text = f"${flat_price:.2f}"
            p = tb.text_frame.paragraphs[0]
            p.alignment = 2
            r = p.runs[0]
            r.font.size = Pt(72)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0x3E, 0x4F, 0xE0)
            r.font.name = "Titillium Web"
            tb2 = slide16.shapes.add_textbox(tb_left, tb_top + Inches(2.2), tb_width, Inches(0.4))
            tb2.text_frame.text = "USD por transacción aprobada"
            p2 = tb2.text_frame.paragraphs[0]
            p2.alignment = 2
            r2 = p2.runs[0]
            r2.font.size = Pt(14)
            r2.font.color.rgb = RGBColor(0x92, 0x95, 0x9B)
            r2.font.name = "Titillium Web"
            tb3 = slide16.shapes.add_textbox(tb_left, tb_top + Inches(2.7), tb_width, Inches(0.35))
            tb3.text_frame.text = "tarifa única sin tramos"
            p3 = tb3.text_frame.paragraphs[0]
            p3.alignment = 2
            r3 = p3.runs[0]
            r3.font.size = Pt(11)
            r3.font.italic = True
            r3.font.color.rgb = RGBColor(0x92, 0x95, 0x9B)
            r3.font.name = "Titillium Web"
            if discount_shape:
                for para in discount_shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""
        elif pricing_type == "tramos" and table_s16 is not None:
            tbl = table_s16.table._tbl
            ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            trs = tbl.findall(f".//{ns}tr")
            current_data_rows = len(trs) - 1
            needed = len(tramos)
            if needed < current_data_rows:
                for i in range(current_data_rows - needed):
                    tr = trs[-(1 + i)]
                    tr.getparent().remove(tr)
            elif needed > current_data_rows:
                last_tr = trs[-1]
                for _ in range(needed - current_data_rows):
                    new_tr = deepcopy(last_tr)
                    last_tr.addnext(new_tr)
            for shape in slide16.shapes:
                if shape.has_table:
                    table_s16 = shape
                    break
            tbl_obj = table_s16.table
            for i, tramo in enumerate(tramos):
                row = tbl_obj.rows[i + 1]
                set_cell_text_preserve(row.cells[0], tramo.get("name", f"Tramo {i+1}"))
                set_cell_text_preserve(row.cells[1], tramo["range"])
                price = tramo["price"]
                price_str = f"${price:.4f}" if price < 0.01 else f"${price:.2f}"
                set_cell_text_preserve(row.cells[2], price_str)
            if discount_shape:
                example_parts = [f"(Trx tramo {i+1} × ${t['price']:.2f})" for i, t in enumerate(tramos[:2])]
                example = " + ".join(example_parts) + " + ..."
                new_text = f"cada tramo se suma al siguiente. Ej: {example}"
                if discount_shape.text_frame.paragraphs and discount_shape.text_frame.paragraphs[0].runs:
                    discount_shape.text_frame.paragraphs[0].runs[0].text = new_text
                    for r in list(discount_shape.text_frame.paragraphs[0].runs[1:]):
                        r_elem = r._r
                        r_elem.getparent().remove(r_elem)

    # SLIDE 17: BDM info
    if len(prs.slides) > 16:
        for shape in prs.slides[16].shapes:
            replace_in_shape(shape, "Rasheed Bayter", bdm_name)
            replace_in_shape(shape, "+57 3164730919", bdm_phone)
            replace_in_shape(shape, "rasheed@y.uno", bdm_email)

    prs.save(output_path)
    return {
        "company": company,
        "tpv_mensual": tpv_mensual,
        "reduccion_mdr": reduccion_mdr,
        "revenue_aumento": revenue_aumento,
        "total_impacto": total_impacto,
    }


def main():
    # Mode detection:
    # - Arg "-" or no args with stdin: read config from stdin, write PPTX to stdout
    # - Args: python generate.py config.json [template.pptx] [output.pptx]

    stdin_mode = (len(sys.argv) == 1) or (len(sys.argv) >= 2 and sys.argv[1] == "-")

    script_dir = Path(__file__).parent

    if stdin_mode:
        config = json.loads(sys.stdin.read())
        template = str(script_dir / "template.pptx")
        # Write to a temp file, then read bytes
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_out:
            output = tmp_out.name
    else:
        config_path = sys.argv[1]
        with open(config_path) as f:
            config = json.load(f)
        template = sys.argv[2] if len(sys.argv) > 2 else str(script_dir / "template.pptx")
        company = config["clientName"]
        default_output = f"Business_Case_{company.replace(' ', '_')}.pptx"
        output = sys.argv[3] if len(sys.argv) > 3 else config.get("outputFile", default_output)

    # Logo: use provided path or generate fallback
    logo_path = config.get("logoPath")
    logo_tmp = None
    if not logo_path:
        color = config.get("brandColor", "3E4FE0")
        company_safe = config["clientName"].replace(" ", "_").replace("/", "_")
        logo_tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False, prefix=f"{company_safe}_logo_")
        logo_tmp.close()
        logo_path = logo_tmp.name
        try:
            create_logo_fallback(config["clientName"], color, logo_path)
        except Exception as e:
            print(f"Logo fallback failed: {e}", file=sys.stderr)
            logo_path = None

    try:
        result = edit_template(config, template, output, logo_path)

        if stdin_mode:
            # Read output PPTX and write bytes to stdout
            with open(output, "rb") as f:
                sys.stdout.buffer.write(f.read())
            # Write summary to stderr as JSON
            print(json.dumps({
                "summary": {
                    "clientName": result["company"],
                    "totalTPVMensual": result["tpv_mensual"],
                    "ahorroMDRMensual": result["reduccion_mdr"],
                    "aumentoRevenue": result["revenue_aumento"],
                    "totalMensual": result["total_impacto"],
                    "slides": 17,
                }
            }), file=sys.stderr)
            os.unlink(output)
        else:
            print(f"Saved: {output}", file=sys.stderr)
            print(f"  TPV: {fmt_money(result['tpv_mensual'])}", file=sys.stderr)
            print(f"  Total impacto: {fmt_k(result['total_impacto'])}/mes", file=sys.stderr)
    finally:
        if logo_tmp and logo_path and os.path.exists(logo_path):
            try:
                os.unlink(logo_path)
            except:
                pass


if __name__ == "__main__":
    main()
